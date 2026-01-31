#!/usr/bin/env python3
"""Generate Splash pitch PPT - Light Theme"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Light Theme Colors
BG_COLOR = RGBColor(0xf8, 0xfa, 0xfc)  # Light gray-blue
PRIMARY = RGBColor(0x4f, 0x46, 0xe5)   # Indigo
ACCENT = RGBColor(0x7c, 0x3a, 0xed)    # Purple
CYAN = RGBColor(0x05, 0x91, 0xb3)      # Cyan darker for contrast
ORANGE = RGBColor(0xea, 0x58, 0x0c)    # Orange
GREEN = RGBColor(0x16, 0xa3, 0x4a)     # Green
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x1e, 0x29, 0x3b)     # Dark slate
DARK_GRAY = RGBColor(0x47, 0x55, 0x69) # Slate gray
MUTED = RGBColor(0x64, 0x74, 0x8b)     # Muted gray
LIGHT_BG = RGBColor(0xf1, 0xf5, 0xf9)  # Lighter gray for cards

# Background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
background.fill.solid()
background.fill.fore_color.rgb = BG_COLOR
background.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=12, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with specified properties"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# === HEADER ===
# Logo
logo = add_rounded_rect(slide, 0.4, 0.3, 0.5, 0.5, PRIMARY)
logo_text = add_text_box(slide, 0.4, 0.35, 0.5, 0.5, "S", 24, WHITE, True, PP_ALIGN.CENTER)

# Title
add_text_box(slide, 1.0, 0.3, 4, 0.4, "Splash", 28, BLACK, True)
add_text_box(slide, 1.0, 0.6, 6, 0.3, "AI 时代的原生 UI 语言 — 让 AI 说人话，让 UI 写自己", 11, MUTED)

# === LEFT: ARCHITECTURE ===
# Section title
add_text_box(slide, 0.4, 1.1, 2, 0.3, "架构层", 14, PRIMARY, True)

# Layer boxes
layer_y = 1.5
layer_height = 0.7
layer_gap = 0.15

# AI Layer (Purple)
ai_bg = add_rounded_rect(slide, 0.4, layer_y, 5.8, layer_height, RGBColor(0xf5, 0xf3, 0xff), ACCENT)
add_text_box(slide, 0.5, layer_y + 0.05, 2, 0.25, "AI Agent Layer", 10, ACCENT, True)
add_text_box(slide, 4.5, layer_y + 0.05, 1.5, 0.25, "LLM / Agent", 8, ACCENT)
# AI boxes
for i, txt in enumerate(["Natural Language", "Intent Parser", "Code Generator"]):
    box = add_rounded_rect(slide, 0.6 + i * 1.8, layer_y + 0.35, 1.6, 0.28, ACCENT)
    add_text_box(slide, 0.6 + i * 1.8, layer_y + 0.38, 1.6, 0.25, txt, 8, WHITE, False, PP_ALIGN.CENTER)

layer_y += layer_height + layer_gap

# Splash Layer (Indigo)
splash_bg = add_rounded_rect(slide, 0.4, layer_y, 5.8, layer_height + 0.35, RGBColor(0xee, 0xf2, 0xff), PRIMARY)
add_text_box(slide, 0.5, layer_y + 0.05, 2.5, 0.25, "Splash Script Layer", 10, PRIMARY, True)
add_text_box(slide, 4.5, layer_y + 0.05, 1.5, 0.25, "Splash DSL", 8, PRIMARY)
for i, txt in enumerate(["Script Parser", "Live Compiler", "Hot Reload"]):
    box = add_rounded_rect(slide, 0.6 + i * 1.8, layer_y + 0.35, 1.6, 0.28, PRIMARY)
    add_text_box(slide, 0.6 + i * 1.8, layer_y + 0.38, 1.6, 0.25, txt, 8, WHITE, False, PP_ALIGN.CENTER)
for i, txt in enumerate(["Widget Builder", "Layout Engine", "Event Binding"]):
    box = add_rounded_rect(slide, 0.6 + i * 1.8, layer_y + 0.7, 1.6, 0.28, PRIMARY)
    add_text_box(slide, 0.6 + i * 1.8, layer_y + 0.73, 1.6, 0.25, txt, 8, WHITE, False, PP_ALIGN.CENTER)

layer_y += layer_height + 0.35 + layer_gap

# Makepad Layer (Cyan)
makepad_bg = add_rounded_rect(slide, 0.4, layer_y, 5.8, layer_height + 0.35, RGBColor(0xec, 0xfe, 0xff), CYAN)
add_text_box(slide, 0.5, layer_y + 0.05, 2.5, 0.25, "Makepad Framework", 10, CYAN, True)
add_text_box(slide, 4.5, layer_y + 0.05, 1.5, 0.25, "Rust", 8, CYAN)
for i, txt in enumerate(["Widgets", "Rendering", "Animation"]):
    box = add_rounded_rect(slide, 0.6 + i * 1.8, layer_y + 0.35, 1.6, 0.28, CYAN)
    add_text_box(slide, 0.6 + i * 1.8, layer_y + 0.38, 1.6, 0.25, txt, 8, WHITE, False, PP_ALIGN.CENTER)
for i, txt in enumerate(["GPU Shaders", "Layout", "Events"]):
    box = add_rounded_rect(slide, 0.6 + i * 1.8, layer_y + 0.7, 1.6, 0.28, CYAN)
    add_text_box(slide, 0.6 + i * 1.8, layer_y + 0.73, 1.6, 0.25, txt, 8, WHITE, False, PP_ALIGN.CENTER)

layer_y += layer_height + 0.35 + layer_gap

# Platform Layer (Orange)
platform_bg = add_rounded_rect(slide, 0.4, layer_y, 5.8, layer_height, RGBColor(0xff, 0xf7, 0xed), ORANGE)
add_text_box(slide, 0.5, layer_y + 0.05, 2, 0.25, "Platform Layer", 10, ORANGE, True)
add_text_box(slide, 4.5, layer_y + 0.05, 1.5, 0.25, "Native", 8, ORANGE)
platforms = ["macOS", "Windows", "Linux", "iOS", "Android", "WASM"]
for i, txt in enumerate(platforms):
    box = add_rounded_rect(slide, 0.5 + i * 0.9, layer_y + 0.35, 0.8, 0.28, ORANGE)
    add_text_box(slide, 0.5 + i * 0.9, layer_y + 0.38, 0.8, 0.25, txt, 7, WHITE, False, PP_ALIGN.CENTER)

# === RIGHT COLUMN ===
right_x = 6.5

# Scenarios
add_text_box(slide, right_x, 1.1, 2, 0.3, "应用场景", 14, PRIMARY, True)
scenarios = [
    ("🤖", "AI Agent 动态生成界面"),
    ("📊", "实时数据仪表盘"),
    ("💬", "对话式 UI 原型设计"),
    ("📱", "跨平台应用快速开发"),
]
for i, (icon, text) in enumerate(scenarios):
    row = i // 2
    col = i % 2
    x = right_x + col * 3.2
    y = 1.45 + row * 0.55
    box = add_rounded_rect(slide, x, y, 3.0, 0.45, WHITE, RGBColor(0xe2, 0xe8, 0xf0))
    add_text_box(slide, x + 0.1, y + 0.1, 2.8, 0.3, f"{icon} {text}", 10, DARK_GRAY)

# Core Tech
add_text_box(slide, right_x, 2.6, 2, 0.3, "核心技术", 14, PRIMARY, True)
techs = [
    ("⚡", "60fps GPU 渲染", "媲美游戏引擎的流畅动画"),
    ("🦀", "Rust 编译原生", "脚本编译为原生代码"),
    ("🔄", "毫秒级热重载", "代码改动即时生效"),
    ("🌍", "全平台支持", "Desktop / Mobile / Web"),
]
for i, (icon, title, desc) in enumerate(techs):
    row = i // 2
    col = i % 2
    x = right_x + col * 3.2
    y = 2.95 + row * 0.7
    box = add_rounded_rect(slide, x, y, 3.0, 0.6, WHITE, RGBColor(0xe2, 0xe8, 0xf0))
    add_text_box(slide, x + 0.1, y + 0.08, 2.8, 0.25, f"{icon} {title}", 10, BLACK, True)
    add_text_box(slide, x + 0.35, y + 0.32, 2.6, 0.25, desc, 8, MUTED)

# Comparison
add_text_box(slide, right_x, 4.4, 2, 0.3, "竞争优势", 14, PRIMARY, True)
comparisons = [
    ("vs Figma", "生成真正可运行的原生应用"),
    ("vs Flutter", "AI-First 架构，动态生成"),
    ("vs React Native", "纯 Rust 栈，无 JS 桥接"),
    ("vs 传统开发", "自然语言 → UI，10x 效率"),
]
for i, (vs, advantage) in enumerate(comparisons):
    y = 4.75 + i * 0.42
    box = add_rounded_rect(slide, right_x, y, 6.4, 0.36, WHITE, RGBColor(0xe2, 0xe8, 0xf0))
    # VS label
    vs_box = add_rounded_rect(slide, right_x + 0.1, y + 0.06, 1.1, 0.24, RGBColor(0xff, 0xed, 0xd5))
    add_text_box(slide, right_x + 0.1, y + 0.08, 1.1, 0.24, vs, 8, ORANGE, True, PP_ALIGN.CENTER)
    add_text_box(slide, right_x + 1.35, y + 0.08, 4.9, 0.24, advantage, 9, PRIMARY)

# === FEATURE HIGHLIGHTS ===
feature_y = 6.15
feature_box = add_rounded_rect(slide, 0.4, feature_y, 12.5, 1.0, WHITE, RGBColor(0xc7, 0xd2, 0xfe))

# Feature intro text
intro = "Splash 是专为 AI 时代设计的下一代 UI 开发语言。将自然语言理解与原生性能渲染深度融合，让开发者和 AI Agent 通过简单指令实时生成高性能跨平台界面。基于 Makepad 框架和 Rust 语言，在保持脚本灵活性的同时实现媲美 C++ 的运行时性能。"
add_text_box(slide, 0.5, feature_y + 0.08, 12.3, 0.5, intro, 9, DARK_GRAY)

# Three feature points
features = [
    ("🤖 AI-First", "原生支持 LLM 输出，流式生成渲染", ACCENT),
    ("⚡ 原生性能", "GPU 加速 60fps，无桥接开销", CYAN),
    ("🚀 极致体验", "毫秒热重载，全平台单代码库", GREEN),
]
for i, (title, desc, color) in enumerate(features):
    x = 0.6 + i * 4.15
    add_text_box(slide, x, feature_y + 0.55, 1.5, 0.2, title, 9, color, True)
    add_text_box(slide, x, feature_y + 0.75, 3.8, 0.2, desc, 8, MUTED)

# === FOOTER ===
add_text_box(slide, 0.4, 7.2, 6, 0.25, '"预测未来的最好方式是创造它。" — Alan Kay', 10, MUTED)
add_text_box(slide, 9.5, 7.2, 3.5, 0.25, "github.com/makepad/makepad  |  robius.rs", 9, MUTED, False, PP_ALIGN.RIGHT)

# Save
output_path = "splash-pitch-light.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
